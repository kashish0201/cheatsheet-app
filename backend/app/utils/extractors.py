from pptx import Presentation
import pdfplumber

def extract_text_from_pptx(file_path:str) -> str:
    text = []
    presentation = Presentation(file_path)

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape,"text"):
                text.append(shape.text)

    return "\n".join(text)

def extract_text_from_pdf(file_path: str) -> str:
    text = []

    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text.append(page_text)

    return "\n".join(text)